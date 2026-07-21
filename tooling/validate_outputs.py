#!/usr/bin/env python3
"""Validate generated deliverables after the build step.

The checks are structural rather than editorial: the script proves that every
declared format exists, opens, contains substantive content and remains
consistent with the source manifest. It writes a machine-readable QA record to
``dist/output-qa.json`` and fails closed on missing or corrupt deliverables.
"""

from __future__ import annotations

import json
import hashlib
import re
import sys
from xml.etree import ElementTree
from pathlib import Path
from zipfile import BadZipFile, ZipFile

import yaml
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
from pypdf import PdfReader


ROOT = Path(__file__).resolve().parents[1]
DIST = ROOT / "dist"
MANIFEST_PATH = ROOT / "documents" / "manifest.yaml"


def load_manifest() -> list[dict]:
    return yaml.safe_load(MANIFEST_PATH.read_text(encoding="utf-8"))["artifacts"]


def slide_count(source: Path) -> int:
    text = source.read_text(encoding="utf-8")
    return len([part for part in re.split(r"\n---\s*\n", text) if part.strip()])


def expected_output_path(artifact: dict, extension: str) -> Path:
    directory = DIST / artifact["id"]
    if extension == "html":
        return directory / "index.html"
    return directory / f"{artifact['id']}.{extension}"


def validate_html(path: Path, artifact: dict, errors: list[str]) -> dict:
    text = path.read_text(encoding="utf-8")
    soup = BeautifulSoup(text, "html.parser")
    headings = soup.find_all(["h1", "h2", "h3"])
    if not soup.title or not headings:
        errors.append(f"{artifact['id']}: HTML lacks title or headings")
    for image in soup.find_all("img"):
        src = image.get("src", "")
        if src.startswith(("http://", "https://", "data:")):
            continue
        candidate = (path.parent / src).resolve()
        if not candidate.exists() or DIST not in candidate.parents:
            errors.append(f"{artifact['id']}: missing or unsafe HTML image {src}")
    return {"bytes": path.stat().st_size, "headings": len(headings)}


def validate_catalog(errors: list[str]) -> dict:
    path = DIST / "index.html"
    if not path.exists():
        errors.append("missing dist/index.html")
        return {"links": 0}
    soup = BeautifulSoup(path.read_text(encoding="utf-8"), "html.parser")
    links = [anchor.get("href") for anchor in soup.find_all("a") if anchor.get("href")]
    for href in links:
        if href.startswith(("http://", "https://", "mailto:", "#")):
            continue
        candidate = (path.parent / href).resolve()
        if not candidate.exists() or DIST not in candidate.parents:
            errors.append(f"catalog: missing or unsafe link {href}")
    return {"links": len(links)}


def validate_pdf(path: Path, artifact: dict, errors: list[str]) -> dict:
    page_text_chars: list[int] = []
    try:
        reader = PdfReader(path)
        pages = len(reader.pages)
        page_text_chars = [
            len("".join((page.extract_text() or "").split())) for page in reader.pages
        ]
        root = reader.trailer["/Root"]
        mark_info = root.get("/MarkInfo")
        if hasattr(mark_info, "get_object"):
            mark_info = mark_info.get_object()
        tagged = bool(root.get("/StructTreeRoot")) and bool(
            mark_info and mark_info.get("/Marked")
        )
    except Exception as exc:  # pypdf exposes several parser-specific errors
        errors.append(f"{artifact['id']}: PDF cannot be parsed: {exc}")
        return {"bytes": path.stat().st_size, "pages": 0}
    if pages < 1:
        errors.append(f"{artifact['id']}: PDF has no pages")
    sparse_pages = [
        index for index, char_count in enumerate(page_text_chars, start=1) if char_count < 100
    ]
    if sparse_pages:
        errors.append(
            f"{artifact['id']}: effectively empty or orphaned PDF pages {sparse_pages}"
        )
    if not tagged:
        errors.append(f"{artifact['id']}: PDF is not tagged for accessibility")
    if artifact["kind"] == "deck":
        expected = slide_count(ROOT / artifact["source"])
        if pages != expected:
            errors.append(
                f"{artifact['id']}: PDF has {pages} pages; source has {expected} slides"
            )
    return {
        "bytes": path.stat().st_size,
        "pages": pages,
        "tagged": tagged,
        "minimum_page_text_chars": min(page_text_chars, default=0),
    }


def validate_docx(path: Path, artifact: dict, errors: list[str]) -> dict:
    try:
        document = Document(path)
    except (BadZipFile, KeyError, ValueError) as exc:
        errors.append(f"{artifact['id']}: DOCX cannot be parsed: {exc}")
        return {"bytes": path.stat().st_size, "paragraphs": 0, "tables": 0}
    text = " ".join(paragraph.text for paragraph in document.paragraphs)
    words = len(re.findall(r"\w+", text, flags=re.UNICODE))
    if words < 100:
        errors.append(f"{artifact['id']}: DOCX contains only {words} paragraph words")
    section = document.sections[0]
    width_mm = section.page_width.mm
    height_mm = section.page_height.mm
    is_a4 = abs(width_mm - 210) < 1 and abs(height_mm - 297) < 1
    if not is_a4:
        errors.append(
            f"{artifact['id']}: DOCX is {width_mm:.1f}×{height_mm:.1f} mm, expected A4"
        )
    try:
        with ZipFile(path) as archive:
            document_xml = ElementTree.fromstring(archive.read("word/document.xml"))
        doc_properties = document_xml.findall(
            ".//{http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing}docPr"
        )
    except (BadZipFile, KeyError, ElementTree.ParseError) as exc:
        errors.append(f"{artifact['id']}: DOCX image metadata cannot be parsed: {exc}")
        doc_properties = []
    missing_alt = [
        index
        for index, properties in enumerate(doc_properties, start=1)
        if not (properties.get("descr", "").strip() or properties.get("title", "").strip())
    ]
    if missing_alt:
        errors.append(f"{artifact['id']}: DOCX images without alt text {missing_alt}")
    return {
        "bytes": path.stat().st_size,
        "paragraphs": len(document.paragraphs),
        "tables": len(document.tables),
        "paragraph_words": words,
        "page_size_mm": [round(width_mm, 1), round(height_mm, 1)],
        "images_with_alt": len(doc_properties) - len(missing_alt),
        "images_total": len(doc_properties),
    }


def validate_pptx(path: Path, artifact: dict, errors: list[str]) -> dict:
    try:
        presentation = Presentation(path)
    except (BadZipFile, KeyError, ValueError) as exc:
        errors.append(f"{artifact['id']}: PPTX cannot be parsed: {exc}")
        return {"bytes": path.stat().st_size, "slides": 0}
    expected = slide_count(ROOT / artifact["source"])
    actual = len(presentation.slides)
    if actual != expected:
        errors.append(f"{artifact['id']}: PPTX has {actual} slides; expected {expected}")
    empty_slides = []
    for index, slide in enumerate(presentation.slides, start=1):
        visible_text = " ".join(
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        )
        if len(visible_text) < 20:
            empty_slides.append(index)
    if empty_slides:
        errors.append(f"{artifact['id']}: effectively empty PPTX slides {empty_slides}")
    picture_alt_texts: list[str] = []
    try:
        with ZipFile(path) as archive:
            slide_files = sorted(
                name
                for name in archive.namelist()
                if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)
            )
            for slide_file in slide_files:
                slide_xml = ElementTree.fromstring(archive.read(slide_file))
                for picture in slide_xml.findall(
                    ".//{http://schemas.openxmlformats.org/presentationml/2006/main}pic"
                ):
                    properties = picture.find(
                        ".//{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr"
                    )
                    picture_alt_texts.append(
                        ""
                        if properties is None
                        else (
                            properties.get("descr", "").strip()
                            or properties.get("title", "").strip()
                        )
                    )
    except (BadZipFile, KeyError, ElementTree.ParseError) as exc:
        errors.append(f"{artifact['id']}: PPTX image metadata cannot be parsed: {exc}")
    invalid_alt = [
        index
        for index, alt_text in enumerate(picture_alt_texts, start=1)
        if not alt_text
        or re.fullmatch(r"[^/\\]+\.(?:png|jpe?g|gif|svg)", alt_text, flags=re.I)
    ]
    if invalid_alt:
        errors.append(
            f"{artifact['id']}: PPTX images without semantic alt text {invalid_alt}"
        )
    return {
        "bytes": path.stat().st_size,
        "slides": actual,
        "images_with_semantic_alt": len(picture_alt_texts) - len(invalid_alt),
        "images_total": len(picture_alt_texts),
    }


def validate_workbook(path: Path, errors: list[str]) -> dict:
    expected_sheets = [
        "README",
        "Assumptions",
        "Commitments",
        "Cost model",
        "Funding map",
        "Impact logic",
        "Sensitivity",
        "KPI",
    ]
    try:
        with ZipFile(path) as archive:
            workbook_xml = archive.read("xl/workbook.xml").decode("utf-8")
            sheet_names = re.findall(r'<sheet name="([^"]+)"', workbook_xml)
            worksheet_files = sorted(
                name
                for name in archive.namelist()
                if re.fullmatch(r"xl/worksheets/sheet\d+\.xml", name)
            )
            worksheet_xml = "".join(
                archive.read(name).decode("utf-8") for name in worksheet_files
            )
    except (BadZipFile, KeyError, UnicodeDecodeError) as exc:
        errors.append(f"financing-scenarios: XLSX cannot be parsed: {exc}")
        return {"bytes": path.stat().st_size, "sheets": [], "formulas": 0}
    formula_count = len(re.findall(r"<f(?:\s[^>]*)?>", worksheet_xml))
    if sheet_names != expected_sheets:
        errors.append(
            f"financing-scenarios: unexpected sheets {sheet_names}; expected {expected_sheets}"
        )
    if formula_count < 30:
        errors.append(f"financing-scenarios: only {formula_count} formulas found")
    if any(marker in worksheet_xml for marker in ("#REF!", "#NAME?", "#DIV/0!")):
        errors.append("financing-scenarios: workbook contains a static formula error marker")
    return {
        "bytes": path.stat().st_size,
        "sheets": sheet_names,
        "formulas": formula_count,
    }


def validate_assets(errors: list[str]) -> dict:
    metadata_dir = ROOT / "assets" / "metadata"
    metadata_results = []
    for metadata_path in sorted(metadata_dir.glob("*.yaml")):
        record = yaml.safe_load(metadata_path.read_text(encoding="utf-8")) or {}
        asset_value = record.get("asset")
        if not asset_value:
            errors.append(f"{metadata_path.relative_to(ROOT)}: missing asset path")
            continue
        asset_path = (ROOT / asset_value).resolve()
        if not asset_path.exists() or ROOT not in asset_path.parents:
            errors.append(f"{metadata_path.relative_to(ROOT)}: invalid asset {asset_value}")
            continue
        digest = hashlib.sha256(asset_path.read_bytes()).hexdigest()
        if digest != record.get("sha256"):
            errors.append(
                f"{metadata_path.relative_to(ROOT)}: SHA-256 mismatch for {asset_value}"
            )
        prompt_value = record.get("prompt")
        if prompt_value and not (ROOT / prompt_value).exists():
            errors.append(
                f"{metadata_path.relative_to(ROOT)}: missing prompt {prompt_value}"
            )
        copied_path = DIST / asset_value
        if not copied_path.exists():
            errors.append(f"dist asset copy missing: {copied_path.relative_to(ROOT)}")
        metadata_results.append(
            {"metadata": str(metadata_path.relative_to(ROOT)), "asset": asset_value, "sha256": digest}
        )

    svg_results = []
    for svg_path in sorted((ROOT / "assets" / "generated").glob("*.svg")):
        try:
            root = ElementTree.parse(svg_path).getroot()
        except ElementTree.ParseError as exc:
            errors.append(f"{svg_path.relative_to(ROOT)}: malformed SVG: {exc}")
            continue
        if not root.tag.endswith("svg"):
            errors.append(f"{svg_path.relative_to(ROOT)}: root element is not SVG")
        svg_results.append(str(svg_path.relative_to(ROOT)))
    return {"metadata": metadata_results, "svg": svg_results}


def main() -> int:
    errors: list[str] = []
    results: list[dict] = []
    minimum_sizes = {"html": 500, "pdf": 5_000, "docx": 10_000, "pptx": 10_000}

    if not DIST.exists():
        print("ERROR: dist/ does not exist; run the build first", file=sys.stderr)
        return 1

    artifacts = load_manifest()
    catalog_result = validate_catalog(errors)
    build_manifest_path = DIST / "build-manifest.json"
    if not build_manifest_path.exists():
        errors.append("missing dist/build-manifest.json")
        built_ids: set[str] = set()
    else:
        build_manifest = json.loads(build_manifest_path.read_text(encoding="utf-8"))
        built_ids = {item["id"] for item in build_manifest.get("artifacts", [])}
    declared_ids = {artifact["id"] for artifact in artifacts}
    if built_ids != declared_ids:
        errors.append(
            f"build manifest IDs differ: missing={sorted(declared_ids-built_ids)}, "
            f"extra={sorted(built_ids-declared_ids)}"
        )

    validators = {
        "html": validate_html,
        "pdf": validate_pdf,
        "docx": validate_docx,
        "pptx": validate_pptx,
    }
    for artifact in artifacts:
        record = {"id": artifact["id"], "outputs": {}}
        for extension in artifact["outputs"]:
            path = expected_output_path(artifact, extension)
            if not path.exists():
                errors.append(f"{artifact['id']}: missing {path.relative_to(ROOT)}")
                continue
            if path.stat().st_size < minimum_sizes[extension]:
                errors.append(
                    f"{artifact['id']}: {extension} is implausibly small "
                    f"({path.stat().st_size} bytes)"
                )
            record["outputs"][extension] = validators[extension](path, artifact, errors)
        results.append(record)

    model_path = DIST / "models" / "financing-scenarios.xlsx"
    if not model_path.exists():
        errors.append("missing dist/models/financing-scenarios.xlsx")
        model_result = {}
    else:
        model_result = validate_workbook(model_path, errors)
    asset_result = validate_assets(errors)

    qa = {
        "status": "failed" if errors else "passed",
        "artifact_count": len(artifacts),
        "catalog": catalog_result,
        "artifacts": results,
        "financing_model": model_result,
        "assets": asset_result,
        "errors": errors,
    }
    (DIST / "output-qa.json").write_text(
        json.dumps(qa, ensure_ascii=False, indent=2), encoding="utf-8"
    )

    print(f"Validated outputs for {len(artifacts)} artifacts and one XLSX model.")
    if errors:
        for error in errors:
            print(f"ERROR: {error}", file=sys.stderr)
        return 1
    print("Output QA passed.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
