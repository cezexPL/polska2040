#!/usr/bin/env python3
from __future__ import annotations

import json
import re
import shutil
from datetime import date
from pathlib import Path

import cairosvg
import markdown
import yaml
from bs4 import BeautifulSoup, NavigableString, Tag
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptxInches, Pt as PptxPt
from weasyprint import HTML

from build_model import build_financing_model


ROOT = Path(__file__).resolve().parents[1]
DIST = ROOT / "dist"
TODAY = date.today().isoformat()
DISCLAIMER_PL = "Niezależny projekt analityczny — wersja robocza — nie jest dokumentem rządowym."
DISCLAIMER_EN = "Independent analytical project — working draft — not a government document."


BASE_CSS = r"""
:root { --navy:#0b1f33; --blue:#194f7a; --red:#c7333e; --gold:#c7a24a; --paper:#ffffff; --ink:#17202a; --muted:#5d6873; --line:#d9e0e6; }
* { box-sizing:border-box; }
body { margin:0; color:var(--ink); background:#eef2f5; font-family:"DejaVu Sans",Arial,sans-serif; line-height:1.55; }
main { width:min(980px, calc(100% - 32px)); margin:24px auto; background:var(--paper); padding:64px 72px; box-shadow:0 8px 30px rgba(11,31,51,.12); }
.draft-banner { background:#fff4d7; border-left:5px solid var(--gold); padding:12px 16px; margin-bottom:36px; font-size:.88rem; }
h1,h2,h3 { color:var(--navy); line-height:1.2; page-break-after:avoid; }
h1 { font-size:2.25rem; border-bottom:4px solid var(--red); padding-bottom:16px; }
h2 { margin-top:2.2rem; border-bottom:1px solid var(--line); padding-bottom:8px; }
h3 { color:var(--blue); }
p,li { orphans:3; widows:3; }
a { color:var(--blue); text-decoration-thickness:1px; }
table { width:100%; border-collapse:collapse; margin:20px 0; font-size:.88rem; }
th { background:var(--navy); color:#fff; text-align:left; }
th,td { border:1px solid var(--line); padding:8px 10px; vertical-align:top; }
tr:nth-child(even) td { background:#f7f9fb; }
blockquote { margin:20px 0; padding:10px 18px; border-left:4px solid var(--blue); color:var(--muted); background:#f4f8fb; }
img { max-width:100%; height:auto; }
code { background:#f2f4f6; padding:.1em .3em; }
.metadata { margin-top:48px; color:var(--muted); font-size:.78rem; border-top:1px solid var(--line); padding-top:12px; }
@media(max-width:700px){ main{width:100%;margin:0;padding:32px 20px;box-shadow:none;} h1{font-size:1.8rem;} table{display:block;overflow-x:auto;} }
@page { size:A4; margin:20mm 17mm 20mm 17mm; @bottom-center { content:"Polska 2040 — draft v0.1  •  " counter(page) " / " counter(pages); color:#5d6873; font-size:8pt; } }
@media print { body{background:#fff;} main{width:auto;margin:0;padding:0;box-shadow:none;} a{color:inherit;} }
"""


DECK_CSS = BASE_CSS + r"""
body { background:#111d29; }
main.deck { width:100%; margin:0; padding:0; background:transparent; box-shadow:none; }
.slide { width:13.333in; height:7.5in; margin:18px auto; padding:.72in .85in; background:#fff; position:relative; overflow:hidden; page-break-after:always; box-shadow:0 8px 30px rgba(0,0,0,.25); }
.slide h1 { margin-top:0; font-size:30pt; border-bottom:4px solid var(--red); }
.slide h2 { font-size:20pt; }
.slide p,.slide li { font-size:16pt; line-height:1.35; }
.slide .slide-number { position:absolute; bottom:.25in; right:.45in; font-size:9pt; color:var(--muted); }
.slide .deck-disclaimer { position:absolute; bottom:.25in; left:.45in; font-size:7.5pt; color:var(--muted); }
@page { size:13.333in 7.5in; margin:0; }
@media print { .slide{margin:0;box-shadow:none;} }
"""


def read_manifest() -> list[dict]:
    with (ROOT / "documents" / "manifest.yaml").open(encoding="utf-8") as handle:
        return yaml.safe_load(handle)["artifacts"]


def markdown_to_body(text: str) -> str:
    return markdown.markdown(text, extensions=["extra", "tables", "fenced_code", "footnotes", "sane_lists", "toc"])


def evidence_markdown() -> str:
    sections = []
    configs = [
        ("Źródła", "research/sources.yaml", "sources", ["id", "title", "publisher", "published_at", "tier", "status"]),
        ("Twierdzenia", "research/claims.yaml", "claims", ["id", "type", "status", "confidence", "text", "source_ids"]),
        ("Sprzeczności", "research/contradictions.yaml", "contradictions", ["id", "claim_ids", "status", "summary"]),
        ("Luki danych", "research/gaps.yaml", "gaps", ["id", "question", "impact", "owner", "status"]),
        ("Decyzje", "research/decisions.yaml", "decisions", ["id", "decision", "addressee", "owner", "status"]),
        ("Ryzyka", "research/risks.yaml", "risks", ["id", "risk", "severity", "owner", "status"]),
    ]
    for title, relative_path, key, columns in configs:
        path = ROOT / relative_path
        data = yaml.safe_load(path.read_text(encoding="utf-8")) if path.exists() else {}
        records = (data or {}).get(key, [])
        sections.append(f"\n## {title}\n")
        if not records:
            sections.append("Brak rekordów w bieżącym wydaniu.\n")
            continue
        sections.append("| " + " | ".join(columns) + " |\n")
        sections.append("|" + "|".join(["---"] * len(columns)) + "|\n")
        for record in records:
            values = []
            for column in columns:
                value = record.get(column, "")
                if isinstance(value, list):
                    value = ", ".join(str(item) for item in value)
                values.append(str(value).replace("|", "\\|").replace("\n", " "))
            sections.append("| " + " | ".join(values) + " |\n")
    return "".join(sections)


def strip_frontmatter(text: str) -> str:
    if text.startswith("---\n"):
        parts = text.split("---\n", 2)
        if len(parts) == 3:
            return parts[2]
    return text


def report_html(artifact: dict, markdown_text: str) -> str:
    disclaimer = DISCLAIMER_EN if artifact["language"] == "en" else DISCLAIMER_PL
    body = markdown_to_body(strip_frontmatter(markdown_text))
    return f"""<!doctype html>
<html lang="{artifact['language']}">
<head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1"><meta name="color-scheme" content="light"><title>{artifact['title']}</title><style>{BASE_CSS}</style></head>
<body><main><div class="draft-banner"><strong>{disclaimer}</strong><br>Stan wiedzy / evidence cut-off: {TODAY}</div>{body}<div class="metadata">Identyfikator: {artifact['id']} · status: {artifact['status']} · wygenerowano: {TODAY}</div></main></body></html>"""


def split_slides(text: str) -> list[str]:
    return [part.strip() for part in re.split(r"\n---\s*\n", strip_frontmatter(text)) if part.strip()]


def deck_html(artifact: dict, markdown_text: str) -> str:
    disclaimer = DISCLAIMER_EN if artifact["language"] == "en" else DISCLAIMER_PL
    rendered = []
    slides = split_slides(markdown_text)
    for index, slide in enumerate(slides, start=1):
        rendered.append(f'<section class="slide">{markdown_to_body(slide)}<div class="deck-disclaimer">{disclaimer}</div><div class="slide-number">{index} / {len(slides)}</div></section>')
    return f"""<!doctype html><html lang="{artifact['language']}"><head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1"><title>{artifact['title']}</title><style>{DECK_CSS}</style></head><body><main class="deck">{''.join(rendered)}</main></body></html>"""


def add_inline_text(paragraph, element: Tag):
    for child in element.children:
        if isinstance(child, NavigableString):
            paragraph.add_run(str(child))
        elif isinstance(child, Tag):
            run = paragraph.add_run(child.get_text())
            if child.name in {"strong", "b"}:
                run.bold = True
            if child.name in {"em", "i"}:
                run.italic = True


def html_to_docx(html_text: str, destination: Path, title: str, disclaimer: str):
    soup = BeautifulSoup(html_text, "html.parser")
    document = Document()
    styles = document.styles
    styles["Normal"].font.name = "Aptos"
    styles["Normal"].font.size = Pt(10.5)
    section = document.sections[0]
    section.top_margin = Inches(.7)
    section.bottom_margin = Inches(.7)
    section.left_margin = Inches(.8)
    section.right_margin = Inches(.8)

    heading = document.add_heading(title, 0)
    heading.alignment = WD_ALIGN_PARAGRAPH.LEFT
    notice = document.add_paragraph()
    notice_run = notice.add_run(disclaimer)
    notice_run.bold = True
    notice_run.font.color.rgb = RGBColor(167, 92, 0)

    for element in soup.body.find_all(["h1", "h2", "h3", "h4", "p", "ul", "ol", "table", "blockquote", "img"], recursive=True):
        if element.find_parent(["ul", "ol", "table", "blockquote"]) and element.name not in {"ul", "ol", "table", "blockquote"}:
            continue
        if element.name in {"h1", "h2", "h3", "h4"}:
            level = min(int(element.name[1]), 3)
            document.add_heading(element.get_text(" ", strip=True), level=level)
        elif element.name == "p":
            paragraph = document.add_paragraph()
            add_inline_text(paragraph, element)
        elif element.name in {"ul", "ol"}:
            style = "List Bullet" if element.name == "ul" else "List Number"
            for item in element.find_all("li", recursive=False):
                paragraph = document.add_paragraph(style=style)
                add_inline_text(paragraph, item)
        elif element.name == "blockquote":
            paragraph = document.add_paragraph(style="Quote")
            paragraph.add_run(element.get_text(" ", strip=True))
        elif element.name == "table":
            rows = element.find_all("tr")
            if not rows:
                continue
            col_count = max(len(row.find_all(["th", "td"], recursive=False)) for row in rows)
            table = document.add_table(rows=0, cols=col_count)
            table.style = "Table Grid"
            for row in rows:
                cells = table.add_row().cells
                for idx, cell in enumerate(row.find_all(["th", "td"], recursive=False)):
                    cells[idx].text = cell.get_text(" ", strip=True)
        elif element.name == "img":
            src = element.get("src")
            if src:
                path = (ROOT / src).resolve()
                if path.exists() and ROOT in path.parents:
                    picture_path = path
                    if path.suffix.lower() == ".svg":
                        converted_dir = ROOT / "build" / "docx-media"
                        converted_dir.mkdir(parents=True, exist_ok=True)
                        picture_path = converted_dir / f"{path.stem}.png"
                        cairosvg.svg2png(url=str(path), write_to=str(picture_path), output_width=1600)
                    document.add_picture(str(picture_path), width=Inches(6.2))

    footer = section.footer.paragraphs[0]
    footer.text = f"Polska 2040 · draft v0.1 · {TODAY}"
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    destination.parent.mkdir(parents=True, exist_ok=True)
    document.save(destination)


def text_from_markdown_line(line: str) -> str:
    return re.sub(r"[*_`#]", "", re.sub(r"\[([^]]+)\]\([^)]+\)", r"\1", line)).strip()


def markdown_to_pptx(artifact: dict, text: str, destination: Path):
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)
    disclaimer = DISCLAIMER_EN if artifact["language"] == "en" else DISCLAIMER_PL

    for slide_index, slide_text in enumerate(split_slides(text), start=1):
        lines = [line.rstrip() for line in slide_text.splitlines() if line.strip()]
        title_line = next((line for line in lines if line.startswith("#")), artifact["title"])
        title = text_from_markdown_line(title_line)
        body_lines = [line for line in lines if line is not title_line and not line.lower().startswith("notes:")]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = PptxRGBColor(255, 255, 255)

        title_box = slide.shapes.add_textbox(PptxInches(.75), PptxInches(.55), PptxInches(11.8), PptxInches(1.0))
        title_frame = title_box.text_frame
        title_frame.clear()
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = PptxPt(28)
        title_p.font.bold = True
        title_p.font.color.rgb = PptxRGBColor(11, 31, 51)

        accent = slide.shapes.add_shape(1, PptxInches(.75), PptxInches(1.45), PptxInches(1.7), PptxInches(.06))
        accent.fill.solid()
        accent.fill.fore_color.rgb = PptxRGBColor(199, 51, 62)
        accent.line.fill.background()

        body_box = slide.shapes.add_textbox(PptxInches(.9), PptxInches(1.8), PptxInches(11.55), PptxInches(4.85))
        body_frame = body_box.text_frame
        body_frame.word_wrap = True
        body_frame.clear()
        first = True
        for line in body_lines:
            clean = text_from_markdown_line(line.lstrip("- "))
            if not clean:
                continue
            paragraph = body_frame.paragraphs[0] if first else body_frame.add_paragraph()
            first = False
            paragraph.text = clean
            paragraph.font.size = PptxPt(18 if line.startswith("-") else 15)
            paragraph.font.color.rgb = PptxRGBColor(23, 32, 42)
            paragraph.level = 0
            paragraph.space_after = PptxPt(10)
            if line.startswith("-"):
                paragraph.text = f"• {clean}"

        footer = slide.shapes.add_textbox(PptxInches(.55), PptxInches(7.08), PptxInches(12.2), PptxInches(.22))
        footer_p = footer.text_frame.paragraphs[0]
        footer_p.text = f"{disclaimer}     {slide_index}"
        footer_p.font.size = PptxPt(7.5)
        footer_p.font.color.rgb = PptxRGBColor(93, 104, 115)
        footer_p.alignment = PP_ALIGN.RIGHT

    destination.parent.mkdir(parents=True, exist_ok=True)
    prs.save(destination)


def build_artifact(artifact: dict) -> dict:
    source_path = ROOT / artifact["source"]
    text = source_path.read_text(encoding="utf-8")
    if artifact["id"] == "evidence-book":
        text += evidence_markdown()
    is_deck = artifact["kind"] == "deck"
    html_text = deck_html(artifact, text) if is_deck else report_html(artifact, text)
    output_dir = DIST / artifact["id"]
    output_dir.mkdir(parents=True, exist_ok=True)
    built = []

    if "html" in artifact["outputs"]:
        target = output_dir / "index.html"
        target.write_text(html_text, encoding="utf-8")
        built.append(str(target.relative_to(ROOT)))
    if "pdf" in artifact["outputs"]:
        target = output_dir / f"{artifact['id']}.pdf"
        HTML(string=html_text, base_url=str(ROOT)).write_pdf(target)
        built.append(str(target.relative_to(ROOT)))
    if "docx" in artifact["outputs"]:
        target = output_dir / f"{artifact['id']}.docx"
        body_html = markdown_to_body(strip_frontmatter(text))
        disclaimer = DISCLAIMER_EN if artifact["language"] == "en" else DISCLAIMER_PL
        html_to_docx(f"<html><body>{body_html}</body></html>", target, artifact["title"], disclaimer)
        built.append(str(target.relative_to(ROOT)))
    if "pptx" in artifact["outputs"]:
        target = output_dir / f"{artifact['id']}.pptx"
        markdown_to_pptx(artifact, text, target)
        built.append(str(target.relative_to(ROOT)))
    return {"id": artifact["id"], "title": artifact["title"], "status": artifact["status"], "outputs": built}


def build_catalog(results: list[dict]):
    items = []
    for result in results:
        links = " · ".join(f'<a href="{Path(path).relative_to("dist")}">{Path(path).suffix.lstrip(".").upper() or "HTML"}</a>' for path in result["outputs"])
        items.append(f"<tr><td>{result['title']}</td><td>{result['status']}</td><td>{links}</td></tr>")
    html_text = f"""<!doctype html><html lang="pl"><head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1"><title>Polska 2040 — artefakty</title><style>{BASE_CSS}</style></head><body><main><h1>Polska 2040 — artefakty v0.1</h1><div class="draft-banner"><strong>{DISCLAIMER_PL}</strong><br>Wygenerowano: {TODAY}</div><table><thead><tr><th>Dokument</th><th>Status</th><th>Formaty</th></tr></thead><tbody>{''.join(items)}</tbody></table></main></body></html>"""
    (DIST / "index.html").write_text(html_text, encoding="utf-8")


def main():
    if DIST.exists():
        shutil.rmtree(DIST)
    DIST.mkdir(parents=True)
    results = [build_artifact(artifact) for artifact in read_manifest()]
    build_financing_model(DIST / "models" / "financing-scenarios.xlsx")
    build_catalog(results)
    (DIST / "build-manifest.json").write_text(json.dumps({"generated_at": TODAY, "artifacts": results}, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"Built {len(results)} artifacts in {DIST}")


if __name__ == "__main__":
    main()
